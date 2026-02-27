import os
import json
import pickle
from pathlib import Path
from datetime import datetime
from transformers import pipeline
from dotenv import load_dotenv
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_core.documents import Document
import hashlib
import re

# Файл уншигч сангууд
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    print("⚠️ PyPDF2 суугаагүй: pip install PyPDF2")

try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    print("⚠️ python-docx суугаагүй: pip install python-docx")

try:
    import pandas as pd
    CSV_AVAILABLE = True
except ImportError:
    CSV_AVAILABLE = False
    print("⚠️ pandas суугаагүй: pip install pandas")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️ python-pptx суугаагүй: pip install python-pptx")

load_dotenv()

# === CONFIG ===
EMBEDDING_MODEL = "sentence-transformers/all-distilroberta-v1"
INDEX_FOLDER = "faiss_disk_index"
METADATA_FILE = "disk_metadata.pkl"
SUPPORTED_EXTENSIONS = ['.txt', '.pdf', '.docx', '.doc', '.json', '.jsonl', '.csv', '.md', '.pptx', '.ppt']

class AdvancedDiskSearch:
    def __init__(self, search_paths=None):
        self.search_paths = search_paths or ["D:/", "C:/Users"]
        self.embedding = HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL)
        self.db = None
        self.metadata = {}
        self.all_docs = []
        
        # AI сонголт
        print("\n💡 AI хариулт үүсгэх горимыг сонгоно уу:")
        print("   1) Идэвхгүй (шууд агуулга харуулах)")
        print("   2) Суурь AI (flan-t5-base, хурдан)")
        print("   3) Дэвшилтэт AI (flan-t5-large, илүү сайн)")
        ai_choice = input("\n   Сонголт (1/2/3, default=1): ").strip() or "1"
        
        self.pipe = None
        self.use_ai = False
        
        if ai_choice in ['2', '3']:
            try:
                model_name = "google/flan-t5-large" if ai_choice == '3' else "google/flan-t5-base"
                print(f"\n🔄 AI model ачаалж байна ({model_name})...")
                if ai_choice == '3':
                    print("   ⚠️  Анх удаа бол 1-2 минут үргэлжилнэ...")
                
                self.pipe = pipeline(
                    "text2text-generation",
                    model=model_name,
                    max_new_tokens=200,
                    device=-1
                )
                self.use_ai = True
                print("✅ AI model бэлэн боллоо\n")
            except Exception as e:
                print(f"⚠️ AI model алдаа: {e}")
                print("   Файлын агуулгыг шууд харуулах горимд шилжлээ\n")
        else:
            print("✅ Файлын агуулгыг шууд харуулах горим\n")

    def get_file_hash(self, filepath):
        try:
            with open(filepath, 'rb') as f:
                return hashlib.md5(f.read()).hexdigest()
        except:
            return None

    def should_skip_directory(self, dirpath):
        skip_dirs = {
            'node_modules', '__pycache__', '.git', '.venv', 'venv',
            'AppData', 'Windows', 'Program Files', 'System32',
            '$RECYCLE.BIN', 'Recovery', 'ProgramData',
            'Microsoft VS Code', 'Visual Studio', 'extensions',
            'resources', 'locales', 'vendor', 'build', 'dist'
        }
        return any(skip in dirpath for skip in skip_dirs)

    def read_txt_file(self, filepath):
        encodings = ['utf-8', 'utf-16', 'cp1252', 'latin-1']
        for encoding in encodings:
            try:
                with open(filepath, 'r', encoding=encoding) as f:
                    return f.read()
            except:
                continue
        return None

    def read_pdf_file(self, filepath):
        if not PDF_AVAILABLE:
            return None
        try:
            text = ""
            with open(filepath, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                for page in pdf_reader.pages[:50]:  # Limit to 50 pages
                    extracted = page.extract_text()
                    if extracted:
                        text += extracted + "\n"
            return text.strip()
        except Exception as e:
            print(f"⚠️ PDF уншихад алдаа {filepath}: {e}")
            return None

    def read_docx_file(self, filepath):
        if not DOCX_AVAILABLE:
            return None
        try:
            doc = docx.Document(filepath)
            text = "\n".join([para.text for para in doc.paragraphs])
            return text.strip()
        except Exception as e:
            print(f"⚠️ DOCX уншихад алдаа {filepath}: {e}")
            return None

    def read_csv_file(self, filepath):
        if not CSV_AVAILABLE:
            return None
        try:
            df = pd.read_csv(filepath, encoding='utf-8', nrows=1000)
            return df.to_string()
        except Exception as e:
            try:
                df = pd.read_csv(filepath, encoding='latin-1', nrows=1000)
                return df.to_string()
            except:
                print(f"⚠️ CSV уншихад алдаа {filepath}: {e}")
                return None

    def read_json_file(self, filepath):
        try:
            if filepath.endswith('.jsonl'):
                texts = []
                with open(filepath, 'r', encoding='utf-8') as f:
                    for i, line in enumerate(f):
                        if i >= 100:
                            break
                        data = json.loads(line)
                        texts.append(json.dumps(data, indent=2, ensure_ascii=False))
                return "\n---\n".join(texts)
            else:
                with open(filepath, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    return json.dumps(data, indent=2, ensure_ascii=False)
        except Exception as e:
            print(f"⚠️ JSON уншихад алдаа {filepath}: {e}")
            return None

    def read_pptx_file(self, filepath):
        if not PPTX_AVAILABLE:
            return None
        try:
            prs = Presentation(filepath)
            text = []
            for slide in prs.slides[:50]:  # Limit to 50 slides
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text.append(shape.text)
            return "\n".join(text)
        except Exception as e:
            print(f"⚠️ PPTX уншихад алдаа {filepath}: {e}")
            return None

    def read_file(self, filepath):
        ext = Path(filepath).suffix.lower()
        if ext in ['.txt', '.md', '.log']:
            return self.read_txt_file(filepath)
        elif ext == '.pdf':
            return self.read_pdf_file(filepath)
        elif ext in ['.docx', '.doc']:
            return self.read_docx_file(filepath)
        elif ext == '.csv':
            return self.read_csv_file(filepath)
        elif ext in ['.json', '.jsonl']:
            return self.read_json_file(filepath)
        elif ext in ['.pptx', '.ppt']:
            return self.read_pptx_file(filepath)
        else:
            return None

    def scan_disk(self, max_files=1000, max_size_mb=10):
        print(f"🔍 Диск хайж эхэлж байна: {self.search_paths}")
        print(f"📂 Дэмжигдсэн файлын төрөл: {', '.join(SUPPORTED_EXTENSIONS)}")
        documents = []
        file_count = 0
        max_size_bytes = max_size_mb * 1024 * 1024
        for search_path in self.search_paths:
            if not os.path.exists(search_path):
                print(f"⚠️ Директори олдсонгүй: {search_path}")
                continue
            print(f"\n📁 Хайж байна: {search_path}")
            for root, dirs, files in os.walk(search_path):
                if self.should_skip_directory(root):
                    dirs[:] = []
                    continue
                for filename in files:
                    if file_count >= max_files:
                        print(f"\n⚠️ Хязгаарт хүрлээ: {max_files} файл")
                        break
                    filepath = os.path.join(root, filename)
                    ext = Path(filename).suffix.lower()
                    if ext not in SUPPORTED_EXTENSIONS:
                        continue
                    try:
                        file_size = os.path.getsize(filepath)
                        if file_size > max_size_bytes or file_size == 0:
                            continue
                    except:
                        continue
                    content = self.read_file(filepath)
                    if not content or len(content.strip()) < 50:
                        continue
                    file_hash = self.get_file_hash(filepath)
                    modified_time = datetime.fromtimestamp(os.path.getmtime(filepath))
                    doc = Document(
                        page_content=content[:5000],
                        metadata={
                            "filename": filename,
                            "filepath": filepath,
                            "extension": ext,
                            "size_kb": file_size / 1024,
                            "modified": modified_time.isoformat(),
                            "hash": file_hash
                        }
                    )
                    documents.append(doc)
                    file_count += 1
                    if file_count % 50 == 0:
                        print(f"✅ {file_count} файл боловсруулсан...")
                if file_count >= max_files:
                    break
        print(f"\n✅ Нийт {len(documents)} файл боловсруулсан")
        self.all_docs = documents
        return documents

    def create_index(self, documents):
        if not documents:
            print("❌ Индекс үүсгэх баримт байхгүй")
            return False
        print(f"\n🔄 FAISS индекс үүсгэж байна ({len(documents)} баримт)...")
        try:
            self.db = FAISS.from_documents(documents, self.embedding)
            os.makedirs(INDEX_FOLDER, exist_ok=True)
            self.db.save_local(INDEX_FOLDER)
            self.metadata = {
                "created": datetime.now().isoformat(),
                "num_documents": len(documents),
                "files": [doc.metadata for doc in documents]
            }
            with open(METADATA_FILE, 'wb') as f:
                pickle.dump(self.metadata, f)
            print(f"✅ Индекс амжилттай үүсгэгдлээ: {INDEX_FOLDER}")
            return True
        except Exception as e:
            print(f"❌ Индекс үүсгэхэд алдаа: {e}")
            return False

    def load_index(self):
        if not os.path.exists(INDEX_FOLDER):
            print("⚠️ Индекс олдсонгүй. Эхлээд scan_disk() дуудна уу.")
            return False
        try:
            print("🔄 Индекс ачаалж байна...")
            self.db = FAISS.load_local(INDEX_FOLDER, self.embedding, allow_dangerous_deserialization=True)
            if os.path.exists(METADATA_FILE):
                with open(METADATA_FILE, 'rb') as f:
                    self.metadata = pickle.load(f)
                print(f"✅ Индекс ачаалагдлаа: {self.metadata.get('num_documents', 0)} баримт")
            else:
                print("⚠️ Metadata олдсонгүй")
            return True
        except Exception as e:
            print(f"❌ Индекс ачаалахад алдаа: {e}")
            return False

    def search_by_keyword(self, keyword):
        keyword = keyword.lower()
        results = []
        if self.metadata and "files" in self.metadata:
            for file_meta in self.metadata["files"]:
                if keyword in file_meta.get("filename", "").lower():
                    results.append(file_meta)
        return results

    def semantic_search(self, query, k=5, score_threshold=2.0):
        if not self.db:
            print("❌ Индекс ачаалаагүй байна")
            return []
        try:
            results = self.db.similarity_search_with_score(query, k=k)
            filtered = [(doc, score) for doc, score in results if score < score_threshold]
            if not filtered and results:
                print(f"⚠️ Threshold-оос давсан, бүх үр дүнг харуулж байна")
                filtered = results
            return filtered
        except Exception as e:
            print(f"❌ Хайлтад алдаа: {e}")
            return []

    def extract_smart_info(self, content, query):
        """
        Ерөнхий зориулалтын ухаалаг мэдээлэл задлах
        Ямар ч төрлийн асуултад хариулах боломжтой
        """
        extracted = []
        query_lower = query.lower()
        query_words = set(re.findall(r'\w+', query_lower))
        
        # === 1. Түлхүүр үг-утга хос олох (Key: Value, Key = Value, Key - Value) ===
        kv_patterns = [
            r'([A-Za-zА-Яа-яёүөЁҮӨ\w\s]{2,30})\s*[:\：=]\s*([^\n\r]{3,100})',
            r'([A-Za-zА-Яа-яёүөЁҮӨ\w\s]{2,30})\s*[-–—]\s*([^\n\r]{3,100})',
        ]
        
        for pattern in kv_patterns:
            matches = re.findall(pattern, content)
            for key, value in matches:
                key_clean = key.strip().lower()
                value_clean = value.strip()
                # Асуултын үгтэй холбоотой эсэхийг шалгах
                key_words = set(re.findall(r'\w+', key_clean))
                if query_words & key_words or any(qw in key_clean for qw in query_words if len(qw) > 2):
                    if len(value_clean) > 2 and len(value_clean) < 200:
                        extracted.append(f"✓ {key.strip()}: {value_clean}")
        
        # === 2. Асуултын түлхүүр үгийн эргэн тойрон дахь контекст ===
        content_lower = content.lower()
        for word in query_words:
            if len(word) < 3:
                continue
            # Түлхүүр үг олох
            for match in re.finditer(re.escape(word), content_lower):
                start = max(0, match.start() - 100)
                end = min(len(content), match.end() + 150)
                context = content[start:end].strip()
                
                # Мөр бүтнээр авах
                lines = context.split('\n')
                relevant_lines = []
                for line in lines:
                    if word in line.lower() and len(line.strip()) > 10:
                        clean_line = line.strip()
                        if clean_line not in [e.split(': ', 1)[-1] if ': ' in e else e for e in extracted]:
                            relevant_lines.append(clean_line)
                
                if relevant_lines and len(extracted) < 5:
                    for line in relevant_lines[:2]:
                        if len(line) < 200:
                            extracted.append(f"→ {line}")
                break  # Зөвхөн эхний тохирлыг авах
        
        # === 3. Тоон утга, хувь, хэмжээ олох ===
        number_context_pattern = r'([A-Za-zА-Яа-яёүөЁҮӨ\w\s]{2,25})\s*[:\：]?\s*(\d+(?:[.,]\d+)?)\s*(%|хувь|percent|USD|₮|¥|\$|€|kg|km|м|cm|mm|GB|MB|TB|件|個|人|年|月|日)?'
        num_matches = re.findall(number_context_pattern, content)
        for label, number, unit in num_matches:
            label_clean = label.strip().lower()
            if any(qw in label_clean for qw in query_words if len(qw) > 2):
                unit_str = unit if unit else ""
                extracted.append(f"✓ {label.strip()}: {number}{unit_str}")
        
        # === 4. Огноо, хугацаа олох ===
        date_patterns = [
            r'(\d{4})\s*[年/-]\s*(\d{1,2})\s*[月/-]?\s*(\d{1,2})?\s*日?',
            r'(\d{1,2})[/.-](\d{1,2})[/.-](\d{2,4})',
        ]
        for pattern in date_patterns:
            date_matches = re.findall(pattern, content)
            if date_matches and any(kw in query_lower for kw in ['when', 'date', 'time', 'хэзээ', 'огноо', 'year', 'month', 'он', 'сар']):
                for dm in date_matches[:3]:
                    date_str = "/".join([d for d in dm if d])
                    if date_str not in str(extracted):
                        extracted.append(f"📅 {date_str}")
        
        # === 5. Email, Phone, URL - ерөнхий ===
        common_patterns = {
            'Email': r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}',
            'Phone': r'[\+]?[\d\s\-\(\)]{8,15}',
            'URL': r'https?://[^\s<>"{}|\\^`\[\]]+',
        }
        for label, pattern in common_patterns.items():
            if any(kw in query_lower for kw in [label.lower(), 'contact', 'холбоо', 'утас', 'имэйл', 'link', 'website']):
                matches = re.findall(pattern, content)
                if matches:
                    extracted.append(f"✓ {label}: {matches[0]}")
        
        # === 6. Жагсаалт, bullet point олох ===
        list_patterns = [
            r'[•·◦▪▫●○]\s*([^\n]{5,100})',
            r'^\s*[-\*]\s+([^\n]{5,100})',
            r'^\s*\d+[.)]\s+([^\n]{5,100})',
        ]
        for pattern in list_patterns:
            list_matches = re.findall(pattern, content, re.MULTILINE)
            relevant_items = [item.strip() for item in list_matches 
                           if any(qw in item.lower() for qw in query_words if len(qw) > 2)]
            if relevant_items and len(extracted) < 8:
                for item in relevant_items[:3]:
                    extracted.append(f"• {item}")
        
        # === 7. Давхардал арилгах ===
        unique_extracted = []
        seen = set()
        for item in extracted:
            # Хялбаршуулсан текстээр давхардал шалгах
            simplified = re.sub(r'[^\w\s]', '', item.lower())
            if simplified not in seen and len(simplified) > 5:
                seen.add(simplified)
                unique_extracted.append(item)
        
        return unique_extracted[:8] if unique_extracted else None

    def generate_answer(self, results, question):
        """Сайжруулсан AI хариулт үүсгэх - олон төрлийн асуултад"""
        if not results:
            return "Олдсон файлуудад хариулт байхгүй."

        # 1. Эхлээд ухаалаг задлалт хийх (pattern matching)
        all_extracted = []
        sources_with_info = []
        
        for doc, score in results[:5]:  # Top 5 баримт шалгах
            extracted = self.extract_smart_info(doc.page_content, question)
            if extracted:
                filename = doc.metadata.get('filename', 'Unknown')
                for info in extracted:
                    all_extracted.append(info)
                    sources_with_info.append(filename)
        
        # Хэрэв pattern-ээр олдвол шууд буцаах
        if all_extracted:
            unique_sources = list(set(sources_with_info))
            result = "\n".join(all_extracted)
            result += f"\n\n📚 Эх: {', '.join(unique_sources)}"
            return result
        
        # 2. Pattern олдохгүй бол AI ашиглах
        if not self.use_ai:
            # AI идэвхгүй бол агуулгын хураангуй буцаах
            summaries = []
            for i, (doc, score) in enumerate(results[:2], 1):
                filename = doc.metadata.get('filename', 'Unknown')
                snippet = doc.page_content[:400].strip()
                summaries.append(f"📄 [{i}] {filename}:\n{snippet}...")
            return "\n\n".join(summaries)
        
        # 3. AI хариулт үүсгэх
        snippets = []
        sources = []
        
        for i, (doc, score) in enumerate(results[:3], 1):
            filename = doc.metadata.get("filename", f"source_{i}")
            content = doc.page_content.replace("\n", " ").strip()
            snippet = content[:350]
            snippets.append(f"[{i}] {filename}: {snippet}")
            sources.append(filename)

        context = "\n\n".join(snippets)
        
        # Товч, тодорхой prompt
        prompt = (
            f"Based on these documents, answer the question directly and concisely.\n"
            f"Extract specific facts, dates, names, or numbers if present.\n\n"
            f"{context}\n\n"
            f"Question: {question}\n"
            f"Answer (be brief and specific):"
        )
        
        # Token шалгалт
        if len(prompt.split()) > 400:
            snippets = [s[:250] for s in snippets[:2]]
            context = "\n".join(snippets)
            prompt = f"Answer based on:\n{context}\n\nQ: {question}\nA:"

        try:
            result = self.pipe(
                prompt,
                max_new_tokens=120,
                do_sample=False,
                num_beams=1,
                early_stopping=True
            )
            
            text = ""
            if isinstance(result, list) and len(result) > 0:
                text = result[0].get("generated_text", "") or result[0].get("summary_text", "")
            
            # Хариулт шалгах
            if not text or len(text.strip()) < 10:
                return f"📋 Файлуудад мэдээлэл байгаа боловч AI задлаж чадсангүй.\n\n📚 Эх: {', '.join(sources)}"
            
            text = text.strip()
            
            # Давталт шалгах (AI алдаа)
            words = text.split()
            if len(words) > 3:
                # Нэг үг хэт их давтагдаж байвал
                word_counts = {}
                for word in words:
                    word_counts[word] = word_counts.get(word, 0) + 1
                max_count = max(word_counts.values())
                if max_count > len(words) / 3:
                    return f"⚠️ AI хариулт алдаатай (давталт илэрсэн)\n\n📚 Эх: {', '.join(sources)}"
            
            # Prompt давталт шалгах
            if "Question:" in text or "Answer:" in text:
                # Prompt-ыг давтаж буцаасан
                parts = text.split("Answer:")
                if len(parts) > 1:
                    text = parts[-1].strip()
            
            # Эх сурвалж нэмэх
            if not any(s in text for s in sources):
                text += f"\n\n📚 Эх: {', '.join(sources)}"
            
            return text
            
        except Exception as e:
            return f"⚠️ AI алдаа: {e}\n\n📚 Эх: {', '.join(sources)}"

    def interactive_search(self):
        print("\n" + "="*60)
        print("🧠 Диск Хайлтын AI Систем Бэлэн")
        print("="*60)
        print("💡 Командууд:")
        print("  - 'stats' - статистик харах")
        print("  - 'rescan' - дахин хайх")
        print("  - 'exit' - гарах")
        print("="*60 + "\n")

        while True:
            user_input = input("🔍 Асуулт: ").strip()
            if not user_input:
                continue
            if user_input.lower() == "exit":
                print("👋 Баяртай!")
                break
            if user_input.lower() == "stats":
                self.show_statistics()
                continue
            if user_input.lower() == "rescan":
                print("\n🔄 Дахин хайж, индекс үүсгэж байна...")
                docs = self.scan_disk(max_files=500)
                self.create_index(docs)
                continue

            # Түлхүүр үгээр хайлт
            keyword_results = self.search_by_keyword(user_input)
            if keyword_results:
                print(f"\n[🔎 Түлхүүр үг] {len(keyword_results)} файл олдлоо:")
                for result in keyword_results[:3]:
                    print(f"  📄 {result['filename']} ({result['extension']})")
                    print(f"     📁 {result['filepath']}")

            # Утгын хайлт
            semantic_results = self.semantic_search(user_input, k=5)
            if not semantic_results:
                print("\n❌ Холбогдох мэдээлэл олдсонгүй.")
                print("💡 Зөвлөмж:")
                print("   - Өөр үг хэллэгээр оролдоно уу")
                print("   - 'stats' командаар ямар файлууд байгааг харна уу")
                print("   - 'rescan' командаар дахин хайж үзнэ үү\n")
                continue

            print(f"\n[📚 Утгын хайлт] {len(semantic_results)} баримт олдлоо:")
            for i, (doc, score) in enumerate(semantic_results, 1):
                snippet = doc.page_content[:200].replace("\n", " ") + "..."
                print(f"\n{i}. 🎯 Оноо: {score:.4f}")
                print(f"   📄 Файл: {doc.metadata.get('filename', 'Unknown')}")
                print(f"   📂 Зам: {doc.metadata.get('filepath', 'Unknown')}")
                print(f"   📏 Хэмжээ: {doc.metadata.get('size_kb', 0):.1f} KB")
                print(f"   📝 Агуулга: {snippet}")

            # AI хариулт
            if self.use_ai:
                print(f"\n🤖 AI хариулт үүсгэж байна...")
                answer = self.generate_answer(semantic_results[:3], user_input)
                print(f"\n💡 AI Хариулт:")
                print("="*60)
                print(answer)
                print("="*60)
            
            # Файлын агуулга үргэлж харуулах
            print(f"\n📋 Дэлгэрэнгүй агуулга:")
            for i, (doc, score) in enumerate(semantic_results[:2], 1):
                content = doc.page_content[:600].strip()
                print(f"\n{'─'*60}")
                print(f"[{i}] {doc.metadata.get('filename')}")
                print(f"{'─'*60}")
                print(content)
                if len(doc.page_content) > 600:
                    print(f"... ({len(doc.page_content)-600} тэмдэгт үлдсэн)")
            
            print("\n" + "-"*60 + "\n")

    def show_statistics(self):
        if not self.metadata:
            print("📊 Статистик байхгүй")
            return
        print("\n" + "="*60)
        print("📊 Системийн Статистик")
        print("="*60)
        print(f"📅 Үүссэн: {self.metadata.get('created', 'Unknown')}")
        print(f"📄 Нийт баримт: {self.metadata.get('num_documents', 0)}")
        if "files" in self.metadata:
            extensions = {}
            total_size = 0
            for file_meta in self.metadata["files"]:
                ext = file_meta.get("extension", "unknown")
                size = file_meta.get("size_kb", 0)
                extensions[ext] = extensions.get(ext, 0) + 1
                total_size += size
            print(f"💾 Нийт хэмжээ: {total_size/1024:.2f} MB")
            print("\n📂 Файлын төрөл:")
            for ext, count in sorted(extensions.items(), key=lambda x: x[1], reverse=True):
                print(f"   {ext}: {count} файл")
            print("\n📋 Баримтын жагсаалт:")
            for i, file_meta in enumerate(self.metadata["files"][:20], 1):
                print(f"   {i}. {file_meta.get('filename', 'Unknown')} ({file_meta.get('extension', 'unknown')})")
            if len(self.metadata["files"]) > 20:
                print(f"   ... болон өөр {len(self.metadata['files']) - 20} файл")
        print("="*60 + "\n")

def main():
    print("🚀 Диск Хайлтын Систем Эхэллээ\n")
    print("📁 Хайх директоруудыг оруулна уу (таслалаар тусгаарлана):")
    print("   Жишээ: D:/Documents, D:/Projects, C:/Users/YourName/Desktop")
    print("   Хоосон орхивол D:/ болон C:/Users хайна")
    user_paths = input("\n📂 Директори: ").strip()
    if user_paths:
        search_paths = [p.strip() for p in user_paths.split(",")]
    else:
        search_paths = ["D:/", "C:/Users"]
    searcher = AdvancedDiskSearch(search_paths=search_paths)
    if os.path.exists(INDEX_FOLDER):
        print(f"\n✅ Хадгалсан индекс олдлоо: {INDEX_FOLDER}")
        choice = input("Ашиглах уу? (y/n): ").strip().lower()
        if choice == 'y':
            if searcher.load_index():
                searcher.interactive_search()
                return
    print("\n🔄 Диск хайж, индекс үүсгэж байна...")
    print("⚠️ Энэ удаан үргэлжлэх боломжтой (5-30 минут)")
    max_files = input("\nХамгийн их файлын тоо (default 1000): ").strip()
    max_files = int(max_files) if max_files.isdigit() else 1000
    docs = searcher.scan_disk(max_files=max_files, max_size_mb=10)
    if docs:
        if searcher.create_index(docs):
            searcher.interactive_search()
    else:
        print("❌ Файл олдсонгүй")

if __name__ == "__main__":
    main() 