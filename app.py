import os
import json
import pickle
import hashlib
import re
from pathlib import Path
from datetime import datetime

from dotenv import load_dotenv

from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_core.documents import Document

from transformers import pipeline

# Readers (optional)
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import pandas as pd
    CSV_AVAILABLE = True
except ImportError:
    CSV_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# ====== UI (PySide6) ======
from PySide6.QtCore import Qt, QThread, Signal
from PySide6.QtWidgets import (
    QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout,
    QPushButton, QLabel, QListWidget, QTextEdit, QFileDialog,
    QLineEdit, QProgressBar, QComboBox, QSpinBox, QDoubleSpinBox,
    QMessageBox, QSplitter, QListWidgetItem, QGroupBox, QTabWidget,
    QTableWidget, QTableWidgetItem, QHeaderView, QStatusBar
)

load_dotenv()

# === CONFIG ===
EMBEDDING_MODEL = "sentence-transformers/all-distilroberta-v1"
INDEX_FOLDER = "faiss_disk_index"
METADATA_FILE = "disk_metadata.pkl"
SUPPORTED_EXTENSIONS = [
    '.txt', '.pdf', '.docx', '.doc', '.json',
    '.jsonl', '.csv', '.md', '.pptx', '.ppt'
]


class AdvancedDiskSearch:
    def __init__(self, search_paths=None, ai_mode="off"):
        self.search_paths = search_paths or ["D:/", "C:/Users"]
        self.embedding = HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL)
        self.db = None
        self.metadata = {}
        self.all_docs = []

        self.pipe = None
        self.use_ai = False
        if ai_mode in ("base", "large"):
            model_name = (
                "google/flan-t5-large" if ai_mode == "large"
                else "google/flan-t5-base"
            )
            self.pipe = pipeline(
                "text2text-generation",
                model=model_name,
                max_new_tokens=200,
                device=-1,
            )
            self.use_ai = True

    # ── Chunked hashing (том файлд санах ой хамгаалалт) ──
    def get_file_hash(self, filepath, chunk_size=1024 * 1024):
        try:
            md5 = hashlib.md5()
            with open(filepath, "rb") as f:
                while True:
                    chunk = f.read(chunk_size)
                    if not chunk:
                        break
                    md5.update(chunk)
            return md5.hexdigest()
        except Exception:
            return None

    def should_skip_directory(self, dirpath):
        skip_dirs = {
            'node_modules', '__pycache__', '.git', '.venv', 'venv',
            'AppData', 'Windows', 'Program Files', 'System32',
            '$RECYCLE.BIN', 'Recovery', 'ProgramData',
            'Microsoft VS Code', 'Visual Studio', 'extensions',
            'resources', 'locales', 'vendor', 'build', 'dist',
        }
        lower = dirpath.lower()
        return any(s.lower() in lower for s in skip_dirs)

    # ── Файл уншигчид ──
    def read_txt_file(self, filepath, max_chars=200_000):
        encodings = ['utf-8', 'utf-16', 'cp1252', 'latin-1']
        for encoding in encodings:
            try:
                with open(filepath, 'r', encoding=encoding, errors="ignore") as f:
                    return f.read(max_chars)
            except Exception:
                continue
        return None

    def read_pdf_file(self, filepath):
        if not PDF_AVAILABLE:
            return None
        try:
            text = ""
            with open(filepath, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages[:50]:
                    extracted = page.extract_text()
                    if extracted:
                        text += extracted + "\n"
            return text.strip()
        except Exception:
            return None

    def read_docx_file(self, filepath):
        if not DOCX_AVAILABLE:
            return None
        try:
            doc = docx.Document(filepath)
            return "\n".join(p.text for p in doc.paragraphs).strip()
        except Exception:
            return None

    def read_csv_file(self, filepath):
        if not CSV_AVAILABLE:
            return None
        try:
            df = pd.read_csv(filepath, encoding='utf-8', nrows=1000)
            return df.to_string()
        except Exception:
            try:
                df = pd.read_csv(filepath, encoding='latin-1', nrows=1000)
                return df.to_string()
            except Exception:
                return None

    def read_json_file(self, filepath):
        try:
            if filepath.endswith('.jsonl'):
                texts = []
                with open(filepath, 'r', encoding='utf-8', errors="ignore") as f:
                    for i, line in enumerate(f):
                        if i >= 100:
                            break
                        data = json.loads(line)
                        texts.append(json.dumps(data, indent=2, ensure_ascii=False))
                return "\n---\n".join(texts)
            else:
                with open(filepath, 'r', encoding='utf-8', errors="ignore") as f:
                    data = json.load(f)
                return json.dumps(data, indent=2, ensure_ascii=False)
        except Exception:
            return None

    def read_pptx_file(self, filepath):
        if not PPTX_AVAILABLE:
            return None
        try:
            prs = Presentation(filepath)
            text = []
            for slide in prs.slides[:50]:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text.append(shape.text)
            return "\n".join(text)
        except Exception:
            return None

    def read_file(self, filepath):
        ext = Path(filepath).suffix.lower()
        if ext in ('.txt', '.md', '.log'):
            return self.read_txt_file(filepath)
        if ext == '.pdf':
            return self.read_pdf_file(filepath)
        if ext in ('.docx', '.doc'):
            return self.read_docx_file(filepath)
        if ext == '.csv':
            return self.read_csv_file(filepath)
        if ext in ('.json', '.jsonl'):
            return self.read_json_file(filepath)
        if ext in ('.pptx', '.ppt'):
            return self.read_pptx_file(filepath)
        return None

    # ── Диск хайлт ──
    def scan_disk(self, max_files=1000, max_size_mb=10, progress_cb=None):
        documents = []
        file_count = 0
        max_size_bytes = max_size_mb * 1024 * 1024

        for search_path in self.search_paths:
            if not os.path.exists(search_path):
                continue
            for root, dirs, files in os.walk(search_path):
                if self.should_skip_directory(root):
                    dirs[:] = []
                    continue
                for filename in files:
                    if file_count >= max_files:
                        break
                    ext = Path(filename).suffix.lower()
                    if ext not in SUPPORTED_EXTENSIONS:
                        continue
                    filepath = os.path.join(root, filename)
                    try:
                        file_size = os.path.getsize(filepath)
                        if file_size > max_size_bytes or file_size == 0:
                            continue
                    except Exception:
                        continue

                    content = self.read_file(filepath)
                    if not content or len(content.strip()) < 50:
                        continue

                    file_hash = self.get_file_hash(filepath)
                    modified_time = datetime.fromtimestamp(
                        os.path.getmtime(filepath)
                    )

                    doc = Document(
                        page_content=content[:5000],
                        metadata={
                            "filename": filename,
                            "filepath": filepath,
                            "extension": ext,
                            "size_kb": file_size / 1024,
                            "modified": modified_time.isoformat(),
                            "hash": file_hash,
                        },
                    )
                    documents.append(doc)
                    file_count += 1

                    if progress_cb and file_count % 10 == 0:
                        progress_cb(file_count, max_files, filepath)

                if file_count >= max_files:
                    break
        self.all_docs = documents
        return documents

    # ── Индекс ──
    def create_index(self, documents):
        if not documents:
            return False, "Индекс үүсгэх баримт байхгүй."
        try:
            self.db = FAISS.from_documents(documents, self.embedding)
            os.makedirs(INDEX_FOLDER, exist_ok=True)
            self.db.save_local(INDEX_FOLDER)
            self.metadata = {
                "created": datetime.now().isoformat(),
                "num_documents": len(documents),
                "files": [doc.metadata for doc in documents],
            }
            with open(METADATA_FILE, 'wb') as f:
                pickle.dump(self.metadata, f)
            return True, f"Индекс амжилттай үүслээ ({len(documents)} баримт)."
        except Exception as e:
            return False, str(e)

    def load_index(self):
        if not os.path.exists(INDEX_FOLDER):
            return False, "Индекс хавтас олдсонгүй."
        try:
            self.db = FAISS.load_local(
                INDEX_FOLDER, self.embedding,
                allow_dangerous_deserialization=True,
            )
            if os.path.exists(METADATA_FILE):
                with open(METADATA_FILE, 'rb') as f:
                    self.metadata = pickle.load(f)
            num = self.metadata.get('num_documents', 0)
            return True, f"Индекс ачаалагдлаа: {num} баримт."
        except Exception as e:
            return False, str(e)

    def search_by_keyword(self, keyword):
        keyword = keyword.lower()
        results = []
        if self.metadata and "files" in self.metadata:
            for file_meta in self.metadata["files"]:
                if keyword in file_meta.get("filename", "").lower():
                    results.append(file_meta)
        return results

    def semantic_search(self, query, k=5, score_threshold=2.0):
        if not self.db:
            return []
        try:
            results = self.db.similarity_search_with_score(query, k=k)
            filtered = [
                (doc, score) for doc, score in results
                if score < score_threshold
            ]
            return filtered if filtered else results
        except Exception:
            return []

    # ══════════════════════════════════════════════════════════
    #  CLI хувилбараас авсан БҮРЭН extract_smart_info
    #  7 алхамтай: KV хос, контекст, тоон утга, огноо,
    #  email/phone/URL, жагсаалт, давхардал арилгах
    # ══════════════════════════════════════════════════════════
    def extract_smart_info(self, content, query):
        extracted = []
        query_lower = query.lower()
        query_words = set(re.findall(r'\w+', query_lower))

        # 1. Түлхүүр үг-утга хос (Key: Value, Key = Value, Key - Value)
        kv_patterns = [
            r'([A-Za-zА-Яа-яёүөЁҮӨ\w\s]{2,30})\s*[:\：=]\s*([^\n\r]{3,100})',
            r'([A-Za-zА-Яа-яёүөЁҮӨ\w\s]{2,30})\s*[-–—]\s*([^\n\r]{3,100})',
        ]
        for pattern in kv_patterns:
            for key, value in re.findall(pattern, content):
                key_clean = key.strip().lower()
                value_clean = value.strip()
                key_words = set(re.findall(r'\w+', key_clean))
                if (query_words & key_words or
                        any(qw in key_clean for qw in query_words if len(qw) > 2)):
                    if 2 < len(value_clean) < 200:
                        extracted.append(f"✓ {key.strip()}: {value_clean}")

        # 2. Түлхүүр үгийн эргэн тойрон дахь контекст
        content_lower = content.lower()
        for word in query_words:
            if len(word) < 3:
                continue
            for match in re.finditer(re.escape(word), content_lower):
                start = max(0, match.start() - 100)
                end = min(len(content), match.end() + 150)
                context = content[start:end].strip()
                lines = context.split('\n')
                relevant_lines = []
                for line in lines:
                    if word in line.lower() and len(line.strip()) > 10:
                        clean_line = line.strip()
                        existing_vals = [
                            e.split(': ', 1)[-1] if ': ' in e else e
                            for e in extracted
                        ]
                        if clean_line not in existing_vals:
                            relevant_lines.append(clean_line)
                if relevant_lines and len(extracted) < 5:
                    for line in relevant_lines[:2]:
                        if len(line) < 200:
                            extracted.append(f"→ {line}")
                break

        # 3. Тоон утга, хувь, хэмжээ
        number_pattern = (
            r'([A-Za-zА-Яа-яёүөЁҮӨ\w\s]{2,25})\s*[:\：]?\s*'
            r'(\d+(?:[.,]\d+)?)\s*'
            r'(%|хувь|percent|USD|₮|¥|\$|€|kg|km|м|cm|mm|GB|MB|TB|件|個|人|年|月|日)?'
        )
        for label, number, unit in re.findall(number_pattern, content):
            label_clean = label.strip().lower()
            if any(qw in label_clean for qw in query_words if len(qw) > 2):
                unit_str = unit if unit else ""
                extracted.append(f"✓ {label.strip()}: {number}{unit_str}")

        # 4. Огноо, хугацаа
        date_patterns = [
            r'(\d{4})\s*[年/-]\s*(\d{1,2})\s*[月/-]?\s*(\d{1,2})?\s*日?',
            r'(\d{1,2})[/.-](\d{1,2})[/.-](\d{2,4})',
        ]
        date_keywords = {
            'when', 'date', 'time', 'хэзээ', 'огноо',
            'year', 'month', 'он', 'сар',
        }
        if date_keywords & set(query_lower.split()):
            for pattern in date_patterns:
                for dm in re.findall(pattern, content)[:3]:
                    date_str = "/".join(d for d in dm if d)
                    if date_str not in str(extracted):
                        extracted.append(f"📅 {date_str}")

        # 5. Email, Phone, URL
        contact_patterns = {
            'Email': r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}',
            'Phone': r'[\+]?[\d\s\-\(\)]{8,15}',
            'URL': r'https?://[^\s<>"{}|\\^`\[\]]+',
        }
        contact_keywords = {
            'email', 'contact', 'холбоо', 'утас',
            'имэйл', 'link', 'website', 'phone', 'url',
        }
        if contact_keywords & set(query_lower.split()):
            for label, pattern in contact_patterns.items():
                matches = re.findall(pattern, content)
                if matches:
                    extracted.append(f"✓ {label}: {matches[0]}")

        # 6. Жагсаалт / bullet point
        list_patterns = [
            r'[•·◦▪▫●○]\s*([^\n]{5,100})',
            r'^\s*[-\*]\s+([^\n]{5,100})',
            r'^\s*\d+[.)]\s+([^\n]{5,100})',
        ]
        for pattern in list_patterns:
            list_matches = re.findall(pattern, content, re.MULTILINE)
            relevant = [
                item.strip() for item in list_matches
                if any(qw in item.lower() for qw in query_words if len(qw) > 2)
            ]
            if relevant and len(extracted) < 8:
                for item in relevant[:3]:
                    extracted.append(f"• {item}")

        # 7. Давхардал арилгах
        unique = []
        seen = set()
        for item in extracted:
            simplified = re.sub(r'[^\w\s]', '', item.lower())
            if simplified not in seen and len(simplified) > 5:
                seen.add(simplified)
                unique.append(item)

        return unique[:8] if unique else None

    # ══════════════════════════════════════════════════════════
    #  CLI хувилбараас авсан БҮРЭН generate_answer
    #  Token хязгаар, давталт, prompt давталт шалгалттай
    # ══════════════════════════════════════════════════════════
    def generate_answer(self, results, question):
        if not results:
            return "Олдсон файлуудад хариулт байхгүй."

        # 1. Pattern matching (extract_smart_info)
        all_extracted = []
        sources_with_info = []

        for doc, score in results[:5]:
            extracted = self.extract_smart_info(doc.page_content, question)
            if extracted:
                filename = doc.metadata.get('filename', 'Unknown')
                for info in extracted:
                    all_extracted.append(info)
                    sources_with_info.append(filename)

        if all_extracted:
            unique_sources = sorted(set(sources_with_info))
            result_text = "\n".join(all_extracted)
            result_text += f"\n\n📚 Эх: {', '.join(unique_sources)}"
            return result_text

        # 2. AI идэвхгүй бол агуулгын хураангуй
        if not self.use_ai:
            summaries = []
            for i, (doc, score) in enumerate(results[:2], 1):
                fn = doc.metadata.get('filename', 'Unknown')
                snippet = doc.page_content[:400].strip()
                summaries.append(f"📄 [{i}] {fn}:\n{snippet}...")
            return "\n\n".join(summaries)

        # 3. AI хариулт үүсгэх
        snippets = []
        sources = []
        for i, (doc, score) in enumerate(results[:3], 1):
            fn = doc.metadata.get("filename", f"source_{i}")
            sources.append(fn)
            text = doc.page_content.replace("\n", " ").strip()[:350]
            snippets.append(f"[{i}] {fn}: {text}")

        context = "\n\n".join(snippets)

        prompt = (
            "Based on these documents, answer the question directly "
            "and concisely.\nExtract specific facts, dates, names, "
            "or numbers if present.\n\n"
            f"{context}\n\n"
            f"Question: {question}\n"
            f"Answer (be brief and specific):"
        )

        # Token хязгаар шалгалт
        if len(prompt.split()) > 400:
            snippets = [s[:250] for s in snippets[:2]]
            context = "\n".join(snippets)
            prompt = f"Answer based on:\n{context}\n\nQ: {question}\nA:"

        try:
            result = self.pipe(
                prompt,
                max_new_tokens=120,
                do_sample=False,
                num_beams=1,
                early_stopping=True,
            )
            text = ""
            if isinstance(result, list) and result:
                text = (result[0].get("generated_text", "")
                        or result[0].get("summary_text", ""))

            if not text or len(text.strip()) < 10:
                return (
                    "📋 Файлуудад мэдээлэл байгаа боловч AI задлаж "
                    f"чадсангүй.\n\n📚 Эх: {', '.join(sources)}"
                )

            text = text.strip()

            # Давталт шалгах
            words = text.split()
            if len(words) > 3:
                word_counts = {}
                for w in words:
                    word_counts[w] = word_counts.get(w, 0) + 1
                if max(word_counts.values()) > len(words) / 3:
                    return (
                        "⚠️ AI хариулт алдаатай (давталт илэрсэн)\n\n"
                        f"📚 Эх: {', '.join(sources)}"
                    )

            # Prompt давталт шалгах
            if "Question:" in text or "Answer:" in text:
                parts = text.split("Answer:")
                if len(parts) > 1:
                    text = parts[-1].strip()

            if not any(s in text for s in sources):
                text += f"\n\n📚 Эх: {', '.join(sources)}"

            return text

        except Exception as e:
            return f"⚠️ AI алдаа: {e}\n\n📚 Эх: {', '.join(sources)}"

    # ── Дэлгэрэнгүй статистик (CLI хувилбараас) ──
    def get_detailed_stats(self):
        if not self.metadata or "files" not in self.metadata:
            return None
        extensions = {}
        total_size = 0.0
        for fm in self.metadata["files"]:
            ext = fm.get("extension", "unknown")
            sz = fm.get("size_kb", 0)
            extensions[ext] = extensions.get(ext, 0) + 1
            total_size += sz
        return {
            "created": self.metadata.get("created", "Unknown"),
            "num_documents": self.metadata.get("num_documents", 0),
            "total_size_mb": total_size / 1024,
            "extensions": dict(
                sorted(extensions.items(), key=lambda x: x[1], reverse=True)
            ),
            "files": self.metadata["files"],
        }


# ══════════════════════════════════════════════════════════
#  Background thread
# ══════════════════════════════════════════════════════════
class ScanThread(QThread):
    progress = Signal(int, int, str)
    finished_ok = Signal(object)
    failed = Signal(str)

    def __init__(self, engine: AdvancedDiskSearch, max_files: int, max_size_mb: int):
        super().__init__()
        self.engine = engine
        self.max_files = max_files
        self.max_size_mb = max_size_mb

    def run(self):
        try:
            docs = self.engine.scan_disk(
                self.max_files, self.max_size_mb,
                progress_cb=lambda d, t, fp: self.progress.emit(d, t, fp),
            )
            self.finished_ok.emit(docs)
        except Exception as e:
            self.failed.emit(str(e))


# ══════════════════════════════════════════════════════════
#  MainWindow — бүрэн сайжруулсан UI
# ══════════════════════════════════════════════════════════
class MainWindow(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("🧠 Disk Search Desktop  (FAISS + AI)")
        self.resize(1280, 780)

        self.engine = AdvancedDiskSearch(
            search_paths=["D:/", "C:/Users"], ai_mode="off"
        )
        self.last_results = []

        # Status bar
        self.status = QStatusBar()
        self.setStatusBar(self.status)
        self.status.showMessage("Бэлэн. Индекс ачаалах эсвэл Scan+Build дарна уу.")

        root = QWidget()
        self.setCentralWidget(root)
        main = QVBoxLayout(root)

        # ── Дээд хэсэг: тохиргоо ──
        cfg_group = QGroupBox("⚙️ Тохиргоо")
        cfg_layout = QVBoxLayout(cfg_group)

        row1 = QHBoxLayout()
        cfg_layout.addLayout(row1)

        self.paths_edit = QLineEdit("D:/, C:/Users")
        row1.addWidget(QLabel("📁 Paths:"))
        row1.addWidget(self.paths_edit, 1)

        self.btn_pick = QPushButton("📂 Add Folder…")
        self.btn_pick.clicked.connect(self.pick_folder)
        row1.addWidget(self.btn_pick)

        row2 = QHBoxLayout()
        cfg_layout.addLayout(row2)

        self.ai_combo = QComboBox()
        self.ai_combo.addItem("AI Идэвхгүй", "off")
        self.ai_combo.addItem("flan-t5-base (хурдан)", "base")
        self.ai_combo.addItem("flan-t5-large (сайн)", "large")
        self.ai_combo.currentIndexChanged.connect(self.reload_engine)
        row2.addWidget(QLabel("🤖 AI:"))
        row2.addWidget(self.ai_combo)

        self.max_files = QSpinBox()
        self.max_files.setRange(50, 200_000)
        self.max_files.setValue(1000)
        row2.addWidget(QLabel("Max files:"))
        row2.addWidget(self.max_files)

        self.max_size = QSpinBox()
        self.max_size.setRange(1, 200)
        self.max_size.setValue(10)
        row2.addWidget(QLabel("Max MB:"))
        row2.addWidget(self.max_size)

        self.btn_load = QPushButton("📥 Load Index")
        self.btn_load.clicked.connect(self.load_index)
        row2.addWidget(self.btn_load)

        self.btn_scan = QPushButton("🔄 Scan + Build Index")
        self.btn_scan.clicked.connect(self.scan_and_build)
        row2.addWidget(self.btn_scan)

        main.addWidget(cfg_group)

        # ── Progress ──
        prog_row = QHBoxLayout()
        main.addLayout(prog_row)
        self.progress = QProgressBar()
        self.progress.setValue(0)
        prog_row.addWidget(self.progress, 1)
        self.progress_label = QLabel("")
        prog_row.addWidget(self.progress_label)

        # ── Хайлтын мөр ──
        search_group = QGroupBox("🔍 Хайлт")
        search_layout = QHBoxLayout(search_group)

        self.query_edit = QLineEdit()
        self.query_edit.setPlaceholderText(
            "Асуулт бичнэ үү... (Enter = хайх)"
        )
        self.query_edit.returnPressed.connect(self.do_search)
        search_layout.addWidget(self.query_edit, 1)

        self.topk = QSpinBox()
        self.topk.setRange(1, 20)
        self.topk.setValue(5)
        search_layout.addWidget(QLabel("TopK:"))
        search_layout.addWidget(self.topk)

        self.thresh = QDoubleSpinBox()
        self.thresh.setRange(0.1, 10.0)
        self.thresh.setSingleStep(0.1)
        self.thresh.setValue(2.0)
        search_layout.addWidget(QLabel("Threshold:"))
        search_layout.addWidget(self.thresh)

        self.btn_search = QPushButton("🔍 Search")
        self.btn_search.clicked.connect(self.do_search)
        search_layout.addWidget(self.btn_search)

        self.btn_stats = QPushButton("📊 Stats")
        self.btn_stats.clicked.connect(self.show_stats)
        search_layout.addWidget(self.btn_stats)

        main.addWidget(search_group)

        # ── Splitter: зүүн (үр дүн) / баруун (tabs) ──
        splitter = QSplitter(Qt.Horizontal)
        main.addWidget(splitter, 1)

        # Зүүн тал: үр дүнгийн жагсаалт
        left = QWidget()
        left_l = QVBoxLayout(left)
        left_l.setContentsMargins(0, 0, 0, 0)

        self.result_count_label = QLabel("Үр дүн: 0")
        left_l.addWidget(self.result_count_label)

        self.results_list = QListWidget()
        self.results_list.itemSelectionChanged.connect(self.on_select_result)
        left_l.addWidget(self.results_list, 1)

        # Keyword хайлтын жагсаалт
        self.kw_label = QLabel("Түлхүүр үгээр: 0")
        left_l.addWidget(self.kw_label)
        self.kw_list = QListWidget()
        self.kw_list.setMaximumHeight(120)
        left_l.addWidget(self.kw_list)

        splitter.addWidget(left)

        # Баруун тал: tab-ууд
        right = QWidget()
        right_l = QVBoxLayout(right)
        right_l.setContentsMargins(0, 0, 0, 0)

        self.file_info = QLabel("Файл сонгоогүй.")
        self.file_info.setWordWrap(True)
        right_l.addWidget(self.file_info)

        self.tabs = QTabWidget()
        right_l.addWidget(self.tabs, 1)

        # Tab 1: Preview
        self.preview = QTextEdit()
        self.preview.setReadOnly(True)
        self.tabs.addTab(self.preview, "📄 Preview")

        # Tab 2: AI Answer
        self.answer = QTextEdit()
        self.answer.setReadOnly(True)
        self.tabs.addTab(self.answer, "🤖 AI Answer")

        # Tab 3: Stats table
        self.stats_table = QTableWidget()
        self.stats_table.setColumnCount(2)
        self.stats_table.setHorizontalHeaderLabels(["Төрөл", "Тоо"])
        self.stats_table.horizontalHeader().setStretchLastSection(True)
        self.tabs.addTab(self.stats_table, "📊 Stats")

        splitter.addWidget(right)
        splitter.setSizes([380, 900])

    # ── Folder picker ──
    def pick_folder(self):
        folder = QFileDialog.getExistingDirectory(self, "Хавтас сонгох")
        if folder:
            cur = self.paths_edit.text().strip()
            new_path = folder.replace("\\", "/")
            if cur:
                self.paths_edit.setText(f"{cur}, {new_path}")
            else:
                self.paths_edit.setText(new_path)
            self.reload_engine()

    def reload_engine(self):
        paths = [p.strip() for p in self.paths_edit.text().split(",") if p.strip()]
        ai_mode = self.ai_combo.currentData()
        self.status.showMessage(f"Engine дахин ачаалж байна (AI={ai_mode})…")
        self.engine = AdvancedDiskSearch(search_paths=paths, ai_mode=ai_mode)
        self.status.showMessage("Engine бэлэн.", 3000)

    # ── Index ──
    def load_index(self):
        ok, msg = self.engine.load_index()
        if ok:
            self.status.showMessage(msg, 5000)
            self._refresh_stats_tab()
        QMessageBox.information(
            self, "Load Index", msg if ok else f"Алдаа: {msg}"
        )

    def scan_and_build(self):
        self.reload_engine()
        self.progress.setValue(0)
        self.progress_label.setText("Эхэлж байна…")
        self.status.showMessage("Диск хайж байна…")

        self.scan_thread = ScanThread(
            self.engine,
            max_files=self.max_files.value(),
            max_size_mb=self.max_size.value(),
        )
        self.scan_thread.progress.connect(self.on_scan_progress)
        self.scan_thread.finished_ok.connect(self.on_scan_done)
        self.scan_thread.failed.connect(self.on_scan_failed)
        self.btn_scan.setEnabled(False)
        self.scan_thread.start()

    def on_scan_progress(self, done, total, filepath):
        pct = int(min(100, (done / max(1, total)) * 100))
        self.progress.setValue(pct)
        self.progress_label.setText(f"{done}/{total} — {filepath}")

    def on_scan_done(self, docs):
        self.btn_scan.setEnabled(True)
        ok, msg = self.engine.create_index(docs)
        self.status.showMessage(msg, 5000)
        if ok:
            self._refresh_stats_tab()
        QMessageBox.information(
            self, "Build Index", msg if ok else f"Алдаа: {msg}"
        )

    def on_scan_failed(self, err):
        self.btn_scan.setEnabled(True)
        self.status.showMessage(f"Алдаа: {err}", 5000)
        QMessageBox.critical(self, "Scan Error", err)

    # ── Хайлт ──
    def do_search(self):
        q = self.query_edit.text().strip()
        if not q:
            return
        if not self.engine.db:
            QMessageBox.warning(
                self, "Индекс байхгүй",
                "Эхлээд Load Index эсвэл Scan+Build дарна уу.",
            )
            return

        self.results_list.clear()
        self.kw_list.clear()
        self.preview.clear()
        self.answer.clear()
        self.file_info.setText("Хайж байна…")
        self.status.showMessage(f"Хайж байна: «{q}»…")

        # Keyword hits
        kw = self.engine.search_by_keyword(q)
        self.kw_label.setText(f"Түлхүүр үгээр: {len(kw)}")
        for fm in kw[:10]:
            self.kw_list.addItem(
                f"📄 {fm['filename']}  ({fm['extension']})  —  {fm['filepath']}"
            )

        # Semantic
        results = self.engine.semantic_search(
            q, k=self.topk.value(),
            score_threshold=self.thresh.value(),
        )
        self.last_results = results

        if not results and not kw:
            self.file_info.setText("❌ Холбогдох мэдээлэл олдсонгү��.")
            self.result_count_label.setText("Үр дүн: 0")
            self.status.showMessage("Олдсонгүй.", 3000)
            # Зөвлөмж
            self.answer.setPlainText(
                "💡 Зөвлөмж:\n"
                "  • Өөр үг хэллэгээр оролдоно уу\n"
                "  • 📊 Stats товч дарж ямар файл байгааг шалгана уу\n"
                "  • 🔄 Scan+Build дахин хийж үзнэ үү"
            )
            self.tabs.setCurrentWidget(self.answer)
            return

        self.result_count_label.setText(f"Үр дүн: {len(results)}")

        for idx, (doc, score) in enumerate(results):
            fn = doc.metadata.get("filename", "Unknown")
            fp = doc.metadata.get("filepath", "")
            sz = doc.metadata.get("size_kb", 0)
            item = QListWidgetItem(
                f"{idx+1}. {fn}  (score={score:.4f}, {sz:.0f}KB)"
            )
            item.setData(Qt.UserRole, idx)
            item.setToolTip(fp)
            self.results_list.addItem(item)

        self.file_info.setText(
            f"✅ {len(results)} семантик үр дүн  |  "
            f"{len(kw)} түлхүүр үгийн тохирол"
        )

        # AI Answer
        ans = self.engine.generate_answer(results[:5], q)
        self.answer.setPlainText(ans)
        self.tabs.setCurrentWidget(self.answer)
        self.status.showMessage(
            f"Хайлт дууслаа: {len(results)} үр дүн.", 5000
        )

        if self.results_list.count() > 0:
            self.results_list.setCurrentRow(0)

    def on_select_result(self):
        items = self.results_list.selectedItems()
        if not items:
            return
        idx = items[0].data(Qt.UserRole)
        if idx is None or idx >= len(self.last_results):
            return
        doc, score = self.last_results[idx]
        meta = doc.metadata
        self.file_info.setText(
            f"📄 {meta.get('filename')}  |  score = {score:.4f}\n"
            f"📂 {meta.get('filepath')}\n"
            f"📏 {meta.get('size_kb', 0):.1f} KB  |  "
            f"🕒 {meta.get('modified')}"
        )
        self.preview.setPlainText(doc.page_content[:5000])
        self.tabs.setCurrentWidget(self.preview)

    # ── Дэлгэрэнгүй Stats (CLI-с авсан) ──
    def _refresh_stats_tab(self):
        stats = self.engine.get_detailed_stats()
        if not stats:
            return
        exts = stats["extensions"]
        self.stats_table.setRowCount(len(exts) + 2)

        row = 0
        self.stats_table.setItem(row, 0, QTableWidgetItem("📅 Үүссэн"))
        self.stats_table.setItem(row, 1, QTableWidgetItem(stats["created"]))
        row += 1
        self.stats_table.setItem(row, 0, QTableWidgetItem("📄 Нийт баримт"))
        self.stats_table.setItem(
            row, 1, QTableWidgetItem(str(stats["num_documents"]))
        )
        row += 1
        for ext, cnt in exts.items():
            self.stats_table.setItem(row, 0, QTableWidgetItem(ext))
            self.stats_table.setItem(row, 1, QTableWidgetItem(str(cnt)))
            row += 1

        self.stats_table.resizeColumnsToContents()

    def show_stats(self):
        self._refresh_stats_tab()
        stats = self.engine.get_detailed_stats()
        if not stats:
            QMessageBox.information(
                self, "Stats",
                "Metadata байхгүй. Индекс ачаалах эсвэл Build хийнэ үү.",
            )
            return

        msg = (
            f"📅 Үүссэн: {stats['created']}\n"
            f"📄 Нийт баримт: {stats['num_documents']}\n"
            f"💾 Нийт хэмжээ: {stats['total_size_mb']:.2f} MB\n\n"
            f"📂 Файлын төрөл:\n"
        )
        for ext, cnt in stats["extensions"].items():
            msg += f"   {ext}: {cnt} файл\n"

        QMessageBox.information(self, "📊 Системийн Статистик", msg)
        self.tabs.setCurrentWidget(self.stats_table)


def main():
    app = QApplication([])
    w = MainWindow()
    w.show()
    app.exec()


if __name__ == "__main__":
    main()